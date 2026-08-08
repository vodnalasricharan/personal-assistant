from __future__ import annotations

import re
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from config.settings import Settings


def _safe_filename(title: str) -> str:
    slug = re.sub(r"[^a-z0-9_]+", "_", title.lower()).strip("_")
    return slug[:40] or "presentation"


def build_presentation(slides: list[dict[str, Any]], settings: Settings) -> Path:
    """Generate a .pptx file from a list of slide dicts and return the file path."""
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for slide_data in slides:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Title box
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), slide_width - Inches(1), Inches(1.1)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_run = title_para.add_run()
        title_run.text = str(slide_data.get("slide_title", ""))
        title_run.font.bold = True
        title_run.font.size = Pt(28)
        title_run.font.color.rgb = RGBColor(0x1F, 0x23, 0x28)

        # Content box with bullet points
        bullet_points: list[str] = slide_data.get("bullet_points", [])
        content_top = Inches(1.6)
        content_box = slide.shapes.add_textbox(
            Inches(0.5), content_top, slide_width - Inches(1), slide_height - content_top - Inches(0.4)
        )
        cf = content_box.text_frame
        cf.word_wrap = True

        for idx, point in enumerate(bullet_points):
            if idx == 0:
                para = cf.paragraphs[0]
            else:
                para = cf.add_paragraph()
            para.space_before = Pt(6)
            run = para.add_run()
            run.text = f"• {point}"
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x35, 0x3B, 0x44)

        # Notes
        notes_text = slide_data.get("notes") or ""
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = str(notes_text)

        # Thin bottom accent bar
        accent_bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            0, slide_height - Inches(0.08), slide_width, Inches(0.08),
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xD4)
        accent_bar.line.fill.background()

    settings.generated_dir.mkdir(parents=True, exist_ok=True)
    filename = f"presentation_{_safe_filename(slides[0].get('slide_title', 'about') if slides else 'empty')}_{uuid.uuid4().hex[:8]}.pptx"
    output_path = settings.generated_dir / filename
    prs.save(str(output_path))
    return output_path
